from __future__ import annotations

import base64
import csv
from pathlib import Path

import docx  # type: ignore
import openpyxl  # type: ignore
from pptx import Presentation  # type: ignore
from bs4 import BeautifulSoup  # type: ignore

from config import LLMConfig
from services.pdf_parser import parse_pdf_to_markdown
from services.storage import DocumentRecord


TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".log"}


def parse_document(
    record: DocumentRecord,
    page_number: int | None,
    llm_config: LLMConfig,
    sheet_name: str | None = None,
    slide_number: int | None = None,
) -> str:
    if record.is_pdf:
        return parse_pdf_to_markdown(record.stored_path, page_number=page_number)

    text = _extract_text(record, sheet_name=sheet_name, slide_number=slide_number)
    if not text:
        raise ValueError("No text could be extracted from this document.")
    if len(text) > llm_config.max_input_chars:
        text = text[: llm_config.max_input_chars]
    return text


def _extract_text(
    record: DocumentRecord,
    sheet_name: str | None = None,
    slide_number: int | None = None,
) -> str:
    extension = record.extension.lower()
    if extension in TEXT_EXTENSIONS or record.content_type.startswith("text/"):
        if extension == ".csv":
            return _read_csv(record.stored_path)
        return record.stored_path.read_text(encoding="utf-8", errors="ignore")

    if extension == ".docx":
        doc = docx.Document(record.stored_path)
        return "\n".join(paragraph.text for paragraph in doc.paragraphs).strip()

    if extension in {".xlsx", ".xlsm"}:
        return _read_excel(record.stored_path, sheet_name=sheet_name)

    if extension == ".pptx":
        return _read_pptx(record.stored_path, slide_number=slide_number)

    if extension in {".html", ".htm"}:
        html = record.stored_path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        return soup.get_text(separator="\n").strip()

    return _decode_binary(record.stored_path)


def _decode_binary(path: Path) -> str:
    data = path.read_bytes()
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return base64.b64encode(data).decode("ascii")


def _read_csv(path: Path) -> str:
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        rows = list(reader)

    if not rows:
        return ""
    headers = rows[0]
    body = rows[1:]

    if not any(headers):
        headers = [f"Column {i + 1}" for i in range(len(rows[0]))]
        body = rows

    header_line = "| " + " | ".join(headers) + " |"
    divider = "| " + " | ".join("---" for _ in headers) + " |"
    body_lines = ["| " + " | ".join(row) + " |" for row in body]
    return "\n".join([header_line, divider, *body_lines]).strip()


def list_excel_sheets(path: Path) -> list[str]:
    workbook = openpyxl.load_workbook(path, data_only=True)
    return [sheet.title for sheet in workbook.worksheets]


def _read_excel(path: Path, sheet_name: str | None = None) -> str:
    workbook = openpyxl.load_workbook(path, data_only=True)
    sheets = workbook.worksheets
    if sheet_name:
        matches = [sheet for sheet in sheets if sheet.title == sheet_name]
        sheets = matches if matches else sheets[:1]

    lines: list[str] = []
    for sheet in sheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            row_values = ["" if cell is None else str(cell) for cell in row]
            lines.append(" | ".join(row_values).strip())
    return "\n".join(lines).strip()


def list_pptx_slides(path: Path) -> list[int]:
    presentation = Presentation(path)
    return list(range(1, len(presentation.slides) + 1))


def _read_pptx(path: Path, slide_number: int | None = None) -> str:
    presentation = Presentation(path)
    total = len(presentation.slides)
    if total == 0:
        return ""

    slides = presentation.slides
    if slide_number is not None:
        if slide_number < 1 or slide_number > total:
            raise ValueError("Slide number out of range")
        slides = [slides[slide_number - 1]]

    lines: list[str] = []
    for index, slide in enumerate(slides, start=1):
        actual_index = slide_number if slide_number is not None else index
        lines.append(f"# Slide {actual_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines.append(text)
    return "\n".join(lines).strip()
